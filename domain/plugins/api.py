"""
插件管理 API 路由
"""

from datetime import UTC, datetime
import html
import io
import json
import os
from pathlib import Path
from typing import Annotated, Any, List

from fastapi import APIRouter, Depends, File, HTTPException, Query, Request, UploadFile
from fastapi.responses import FileResponse

from domain.audit import AuditAction, audit
from domain.auth import User, get_current_active_user
from domain.permission import require_system_permission
from domain.permission.types import SystemPermission
from .service import PluginService
from .types import (
    PluginInstallResult,
    PluginOut,
)

router = APIRouter(prefix="/api/plugins", tags=["plugins"])


# ========== 安装 ==========


@router.post("/install", response_model=PluginInstallResult)
@audit(action=AuditAction.CREATE, description="安装插件包")
@require_system_permission(SystemPermission.ROLE_MANAGE)
async def install_plugin(
    request: Request,
    current_user: Annotated[User, Depends(get_current_active_user)],
    file: UploadFile = File(...),
):
    """
    安装 .foxpkg 插件包

    上传 .foxpkg 文件进行安装。
    """
    content = await file.read()
    return await PluginService.install_package(content, file.filename or "plugin.foxpkg")


# ========== 插件列表和详情 ==========


@router.get("", response_model=List[PluginOut])
@audit(action=AuditAction.READ, description="获取插件列表")
async def list_plugins(
    request: Request,
    current_user: Annotated[User, Depends(get_current_active_user)],
):
    """获取已安装的插件列表"""
    return await PluginService.list_plugins()


@router.get("/catalog")
async def list_plugin_catalog(request: Request):
    """QiQiClaw 本地应用目录。

    远程应用中心已从前端移除；该接口保留应用发现页的本地扩展点，
    默认返回空列表，避免四个模块主动请求外部服务。
    """
    return {"apps": []}


@router.get("/catalog/{key}")
async def get_plugin_catalog_item(request: Request, key: str):
    """QiQiClaw 本地应用目录详情。"""
    raise HTTPException(status_code=404, detail="Catalog app not found")



def _file_module_data_dir() -> Path:
    base = os.environ.get("QIQICLAW_FILE_MODULE_DATA_DIR")
    return Path(base).expanduser() if base else Path.home() / ".qiqiclaw" / "file-modules"


def _video_data_root() -> Path:
    return _file_module_data_dir() / "data" / ".video"


def _read_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def _file_mtime_iso(path: Path) -> str:
    try:
        return datetime.fromtimestamp(path.stat().st_mtime, tz=UTC).isoformat()
    except FileNotFoundError:
        return ""


def _video_title(payload: dict[str, Any]) -> str:
    detail = (payload.get("tmdb") or {}).get("detail") or {}
    if payload.get("type") == "tv":
        return str(detail.get("name") or detail.get("original_name") or "")
    return str(detail.get("title") or detail.get("original_title") or "")


def _video_year(payload: dict[str, Any]) -> str | None:
    detail = (payload.get("tmdb") or {}).get("detail") or {}
    value = detail.get("first_air_date") if payload.get("type") == "tv" else detail.get("release_date")
    return value[:4] if isinstance(value, str) and value else None


def _video_genres(payload: dict[str, Any]) -> list[str]:
    detail = (payload.get("tmdb") or {}).get("detail") or {}
    return [str(g.get("name")) for g in detail.get("genres") or [] if isinstance(g, dict) and g.get("name")]


def _video_summary(item_id: str, payload: dict[str, Any], mtime_iso: str) -> dict[str, Any]:
    detail = (payload.get("tmdb") or {}).get("detail") or {}
    episodes = payload.get("episodes") or []
    seasons = {e.get("season") for e in episodes if isinstance(e, dict) and e.get("season") is not None}
    return {
        "id": item_id,
        "type": payload.get("type") or "unknown",
        "title": _video_title(payload),
        "year": _video_year(payload),
        "overview": detail.get("overview"),
        "poster_path": detail.get("poster_path"),
        "backdrop_path": detail.get("backdrop_path"),
        "genres": _video_genres(payload),
        "tmdb_id": (payload.get("tmdb") or {}).get("id"),
        "source_path": payload.get("source_path"),
        "scraped_at": payload.get("scraped_at"),
        "updated_at": mtime_iso,
        "episodes_count": len(episodes) if isinstance(episodes, list) else 0,
        "seasons_count": len(seasons),
        "vote_average": detail.get("vote_average"),
        "vote_count": detail.get("vote_count"),
    }


@router.get("/video-library/library")
async def builtin_video_library_list(
    q: str | None = Query(None),
    media_type: str | None = Query(None, alias="type"),
):
    keyword = (q or "").strip().lower()
    type_filter = (media_type or "").strip().lower()
    if type_filter and type_filter not in {"tv", "movie"}:
        raise HTTPException(status_code=400, detail="type must be tv or movie")

    items: list[dict[str, Any]] = []
    for sub in ("tv", "movie"):
        folder = _video_data_root() / sub
        if not folder.exists():
            continue
        for path in folder.glob("*.json"):
            try:
                payload = _read_json(path)
            except Exception:
                continue
            if type_filter and str(payload.get("type") or "").lower() != type_filter:
                continue
            summary = _video_summary(path.stem, payload, _file_mtime_iso(path))
            if keyword and keyword not in f"{summary.get('title') or ''} {summary.get('overview') or ''}".lower():
                continue
            items.append(summary)
    items.sort(key=lambda x: x.get("updated_at") or "", reverse=True)
    return {"code": 0, "data": items}


@router.get("/video-library/library/{item_id}")
async def builtin_video_library_get(item_id: str):
    candidates = [_video_data_root() / "tv" / f"{item_id}.json", _video_data_root() / "movie" / f"{item_id}.json"]
    path = next((p for p in candidates if p.exists()), None)
    if not path:
        raise HTTPException(status_code=404, detail="Item not found")
    payload = _read_json(path)
    payload["id"] = item_id
    payload["updated_at"] = _file_mtime_iso(path)
    return {"code": 0, "data": payload}


_MAX_OFFICE_TEXT = 200_000
_MAX_OFFICE_ROWS = 200
_MAX_OFFICE_COLS = 50


def _safe_html(value: Any) -> str:
    return "" if value is None else html.escape(str(value))


def _office_docx(data: bytes, title: str) -> dict[str, Any]:
    import mammoth

    result = mammoth.convert_to_html(io.BytesIO(data))
    messages = "".join(f"<li>{_safe_html(getattr(msg, 'message', msg))}</li>" for msg in result.messages[:20])
    warning = f"<ul class='qc-office-warnings'>{messages}</ul>" if messages else ""
    return {"title": title, "html": warning + result.value}


def _office_xlsx(data: bytes, title: str) -> dict[str, Any]:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets[:10]:
        parts.append(f"<section class='qc-office-sheet'><h2>{_safe_html(ws.title)}</h2><table>")
        row_count = 0
        for row in ws.iter_rows(max_row=_MAX_OFFICE_ROWS, max_col=_MAX_OFFICE_COLS, values_only=True):
            row_count += 1
            tag = "th" if row_count == 1 else "td"
            cells = "".join(f"<{tag}>{_safe_html(cell)}</{tag}>" for cell in row)
            parts.append(f"<tr>{cells}</tr>")
        parts.append("</table></section>")
    return {"title": title, "html": "".join(parts) or "<p>空工作簿</p>"}


def _office_pptx(data: bytes, title: str) -> dict[str, Any]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    slides: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(_safe_html(shape.text))
        body = "".join(f"<p>{text.replace(chr(10), '<br>')}</p>" for text in texts)
        slides.append(f"<section class='qc-office-slide'><h2>Slide {idx}</h2>{body or '<p></p>'}</section>")
    return {"title": title, "html": "".join(slides) or "<p>空演示文稿</p>"}


@router.get("/office-viewer/preview")
async def builtin_office_preview(path: str = Query(..., min_length=1)):
    from api.response import success
    from domain.virtual_fs.service import VirtualFSService

    normalized = path if path.startswith("/") else f"/{path}"
    name = Path(normalized).name or "document"
    ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""
    if ext in {"doc", "xls", "ppt"}:
        raise HTTPException(
            status_code=415,
            detail="旧版 Office 二进制格式需要 LibreOffice/soffice 转换；当前环境未检测到该依赖。请另存为 docx/xlsx/pptx 后预览。",
        )
    try:
        data = await VirtualFSService.read_file(normalized)
    except Exception as exc:
        raise HTTPException(status_code=404, detail=f"读取文件失败: {exc}") from exc
    if not isinstance(data, (bytes, bytearray)):
        data = bytes(data)
    if len(data) > 25 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="Office 文件超过 25MB，本地预览已拒绝。")
    try:
        if ext == "docx":
            payload = _office_docx(bytes(data), name)
        elif ext == "xlsx":
            payload = _office_xlsx(bytes(data), name)
        elif ext == "pptx":
            payload = _office_pptx(bytes(data), name)
        else:
            raise HTTPException(status_code=415, detail=f"不支持的 Office 格式: {ext or 'unknown'}")
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Office 本地预览失败: {exc}") from exc
    payload["html"] = str(payload.get("html") or "")[:_MAX_OFFICE_TEXT]
    return success(payload)


@router.get("/{key_or_id}", response_model=PluginOut)
@audit(action=AuditAction.READ, description="获取插件详情")
async def get_plugin(
    request: Request,
    key_or_id: str,
    current_user: Annotated[User, Depends(get_current_active_user)],
):
    """获取单个插件详情"""
    return await PluginService.get_plugin(key_or_id)


# ========== 插件管理 ==========


@router.delete("/{key_or_id}")
@audit(action=AuditAction.DELETE, description="卸载插件")
@require_system_permission(SystemPermission.ROLE_MANAGE)
async def delete_plugin(
    request: Request,
    key_or_id: str,
    current_user: Annotated[User, Depends(get_current_active_user)],
):
    """卸载插件"""
    await PluginService.delete(key_or_id)
    return {"code": 0, "msg": "ok"}


# ========== 插件资源 ==========


@router.get("/{key_or_id}/bundle.js")
async def get_bundle(request: Request, key_or_id: str):
    """获取插件前端 bundle"""
    path = await PluginService.get_bundle_path(key_or_id)
    v = (request.query_params.get("v") or "").strip()
    cache_control = "public, max-age=31536000, immutable" if v else "no-cache"
    return FileResponse(
        path,
        media_type="application/javascript",
        headers={"Cache-Control": cache_control},
    )


@router.get("/{key}/assets/{asset_path:path}")
async def get_asset(request: Request, key: str, asset_path: str):
    """获取插件静态资源"""
    path = await PluginService.get_asset_path(key, asset_path)

    # 根据扩展名确定 MIME 类型
    ext = path.suffix.lower()
    media_types = {
        ".js": "application/javascript",
        ".css": "text/css",
        ".json": "application/json",
        ".svg": "image/svg+xml",
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".gif": "image/gif",
        ".webp": "image/webp",
        ".ico": "image/x-icon",
        ".woff": "font/woff",
        ".woff2": "font/woff2",
        ".ttf": "font/ttf",
        ".eot": "application/vnd.ms-fontobject",
        ".html": "text/html",
        ".txt": "text/plain",
        ".md": "text/markdown",
    }
    media_type = media_types.get(ext, "application/octet-stream")

    return FileResponse(
        path,
        media_type=media_type,
        headers={"Cache-Control": "public, max-age=3600"},
    )
