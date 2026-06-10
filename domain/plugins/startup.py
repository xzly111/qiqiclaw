"""
插件启动加载模块

负责在应用启动时加载所有已安装的插件
"""

import json
import logging
import os
import shutil
from pathlib import Path
from typing import TYPE_CHECKING, List, Tuple

from .loader import PluginLoadError, PluginLoader
from .types import PluginManifest

if TYPE_CHECKING:
    from fastapi import FastAPI

logger = logging.getLogger(__name__)

_BUILTIN_FORMAT_PLUGIN_KEYS = (
    "cc.foxel.pdfviewer",
    "cc.foxel.imageviewer",
    "cc.foxel.texteditor",
    "cc.foxel.markdowneditor",
    "cc.foxel.officeviewer",
    "cc.foxel.videolibrary",
)


def _builtin_plugin_source_root() -> Path | None:
    candidates = []
    configured = os.environ.get("QIQICLAW_BUILTIN_PLUGIN_SOURCE_DIR")
    if configured:
        candidates.append(Path(configured).expanduser())
    candidates.append(Path("/home/szd/foxel-deploy/data/plugins"))
    for candidate in candidates:
        if candidate.is_dir():
            return candidate
    return None


def _ignore_runtime_artifacts(_src: str, names: list[str]) -> set[str]:
    return {
        name
        for name in names
        if name == "__pycache__" or name.endswith((".pyc", ".pyo"))
    }



def _office_preview_route_source() -> str:
    return """
import html
import io
from pathlib import Path
from typing import Any

from fastapi import APIRouter, HTTPException, Query

from api.response import success
from domain.virtual_fs.service import VirtualFSService

router = APIRouter()

_MAX_TEXT = 200_000
_MAX_ROWS = 200
_MAX_COLS = 50


def _page(title: str, body: str) -> dict[str, Any]:
    return {"title": title, "html": body}


def _safe(value: Any) -> str:
    if value is None:
        return ""
    return html.escape(str(value))


def _preview_docx(data: bytes, title: str) -> dict[str, Any]:
    import mammoth

    result = mammoth.convert_to_html(io.BytesIO(data))
    messages = "".join(
        f"<li>{_safe(getattr(msg, 'message', msg))}</li>" for msg in result.messages[:20]
    )
    warning = f"<ul class='qc-office-warnings'>{messages}</ul>" if messages else ""
    return _page(title, warning + result.value)


def _preview_xlsx(data: bytes, title: str) -> dict[str, Any]:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets[:10]:
        parts.append(f"<section class='qc-office-sheet'><h2>{_safe(ws.title)}</h2><table>")
        row_count = 0
        for row in ws.iter_rows(max_row=_MAX_ROWS, max_col=_MAX_COLS, values_only=True):
            row_count += 1
            tag = "th" if row_count == 1 else "td"
            cells = "".join(f"<{tag}>{_safe(cell)}</{tag}>" for cell in row)
            parts.append(f"<tr>{cells}</tr>")
        parts.append("</table></section>")
    return _page(title, "".join(parts) or "<p>空工作簿</p>")


def _preview_pptx(data: bytes, title: str) -> dict[str, Any]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    slides: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(_safe(shape.text))
        body = "".join(f"<p>{text.replace(chr(10), '<br>')}</p>" for text in texts)
        slides.append(f"<section class='qc-office-slide'><h2>Slide {idx}</h2>{body or '<p></p>'}</section>")
    return _page(title, "".join(slides) or "<p>空演示文稿</p>")


@router.get("/preview")
async def preview(path: str = Query(..., min_length=1)):
    normalized = path if path.startswith("/") else f"/{path}"
    name = Path(normalized).name or "document"
    ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""

    if ext in {"doc", "xls", "ppt"}:
        raise HTTPException(
            status_code=415,
            detail="旧版 Office 二进制格式需要 LibreOffice/soffice 转换；当前环境未检测到该依赖。请另存为 docx/xlsx/pptx 后预览。",
        )

    try:
        data = await VirtualFSService.read_file(normalized)
    except Exception as exc:
        raise HTTPException(status_code=404, detail=f"读取文件失败: {exc}") from exc

    if not isinstance(data, (bytes, bytearray)):
        data = bytes(data)
    if len(data) > 25 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="Office 文件超过 25MB，本地预览已拒绝。")

    try:
        if ext == "docx":
            payload = _preview_docx(bytes(data), name)
        elif ext == "xlsx":
            payload = _preview_xlsx(bytes(data), name)
        elif ext == "pptx":
            payload = _preview_pptx(bytes(data), name)
        else:
            raise HTTPException(status_code=415, detail=f"不支持的 Office 格式: {ext or 'unknown'}")
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Office 本地预览失败: {exc}") from exc

    payload["html"] = str(payload.get("html") or "")[:_MAX_TEXT]
    return success(payload)
"""


def _office_frontend_source() -> str:
    return """
(function(React, ReactDOM, antd) {
  \"use strict\";
  const { Button, Result, Spin, Typography } = antd;
  const css = `
    .qc-office-viewer{height:100%;overflow:auto;background:#f6f7f9;color:#1f2328;box-sizing:border-box;padding:24px;}
    .qc-office-paper{max-width:1080px;margin:0 auto;background:#fff;border:1px solid #d8dee4;box-shadow:0 8px 28px rgba(31,35,40,.08);padding:32px;min-height:calc(100% - 64px);}
    .qc-office-title{font-size:18px;font-weight:600;margin:0 0 20px;color:#111827;word-break:break-all;}
    .qc-office-content{font-size:14px;line-height:1.72;}
    .qc-office-content table{border-collapse:collapse;width:100%;margin:12px 0 24px;font-size:13px;}
    .qc-office-content th,.qc-office-content td{border:1px solid #d0d7de;padding:6px 8px;vertical-align:top;}
    .qc-office-content th{background:#f3f4f6;font-weight:600;}
    .qc-office-content img{max-width:100%;height:auto;}
    .qc-office-sheet,.qc-office-slide{margin-bottom:28px;}
    .qc-office-slide{border:1px solid #d0d7de;border-radius:6px;padding:18px;background:#fff;}
    .qc-office-warnings{padding:10px 16px;margin:0 0 16px;border:1px solid #f0d98c;background:#fff8db;}
  `;
  const i18n = window.__FOXEL_EXTERNALS__.i18n.create({ zh: {
    \"Preparing document...\": \"正在准备文档...\",
    \"Unable to load document\": \"无法加载文档\",
    \"Close\": \"关闭\",
    \"Local Office Preview\": \"本地 Office 预览\"
  }});
  const { useI18n } = i18n;
  function ensureStyle() {
    if (document.getElementById(\"qc-office-viewer-style\")) return;
    const el = document.createElement(\"style\");
    el.id = \"qc-office-viewer-style\";
    el.textContent = css;
    document.head.appendChild(el);
  }
  function Viewer(props) {
    const { filePath, host } = props;
    const { t } = useI18n();
    const { foxelApi } = window.__FOXEL_EXTERNALS__;
    const [loading, setLoading] = React.useState(true);
    const [error, setError] = React.useState();
    const [payload, setPayload] = React.useState();
    React.useEffect(() => {
      let canceled = false;
      ensureStyle();
      setLoading(true);
      setError(undefined);
      setPayload(undefined);
      foxelApi.request(`/plugins/office-viewer/preview?path=${encodeURIComponent(filePath)}`)
        .then((resp) => {
          if (canceled) return;
          const data = resp && typeof resp === \"object\" && \"data\" in resp ? resp.data : resp;
          setPayload(data || {});
        })
        .catch((err) => {
          if (!canceled) setError(err && err.message ? err.message : String(err || \"Unknown error\"));
        })
        .finally(() => { if (!canceled) setLoading(false); });
      return () => { canceled = true; };
    }, [filePath, foxelApi]);
    if (loading) {
      return React.createElement(\"div\", { style: { height: \"100%\", display: \"flex\", alignItems: \"center\", justifyContent: \"center\" } },
        React.createElement(Spin, { tip: t(\"Preparing document...\") })
      );
    }
    if (error) {
      return React.createElement(Result, { status: \"error\", title: t(\"Unable to load document\"), subTitle: error,
        extra: React.createElement(Button, { type: \"primary\", onClick: host.close }, t(\"Close\")) });
    }
    return React.createElement(\"div\", { className: \"qc-office-viewer\" },
      React.createElement(\"article\", { className: \"qc-office-paper\" },
        React.createElement(Typography.Title, { level: 4, className: \"qc-office-title\" }, payload.title || t(\"Local Office Preview\")),
        React.createElement(\"div\", { className: \"qc-office-content\", dangerouslySetInnerHTML: { __html: payload.html || \"\" } })
      )
    );
  }
  if (!window.__FOXEL_EXTERNALS__) throw new Error(\"Foxel externals not found\");
  window.FoxelRegister({ mount: (container, props) => {
    const root = ReactDOM.createRoot(container);
    root.render(React.createElement(Viewer, props));
    return () => root.unmount();
  }});
})(window.__FOXEL_EXTERNALS__.React, window.__FOXEL_EXTERNALS__.ReactDOM, window.__FOXEL_EXTERNALS__.antd);
"""


def _manifest_to_record(rec, manifest: PluginManifest) -> None:
    rec.name = manifest.name
    rec.version = manifest.version
    rec.description = manifest.description
    rec.author = manifest.author
    rec.website = manifest.website
    rec.github = manifest.github
    rec.license = manifest.license
    rec.manifest = manifest.model_dump(mode="json")

    frontend = manifest.frontend
    rec.open_app = bool(frontend.open_app) if frontend else False
    rec.supported_exts = frontend.supported_exts if frontend else None
    rec.default_bounds = frontend.default_bounds if frontend else None
    rec.default_maximized = frontend.default_maximized if frontend else None
    rec.icon = frontend.icon if frontend else None


async def bootstrap_builtin_format_plugins() -> Tuple[int, List[str]]:
    """Install bundled Foxel format viewers into the QiQiClaw plugin store."""
    from models.database import Plugin

    source_root = _builtin_plugin_source_root()
    if source_root is None:
        return 0, ["未找到原项目内置格式插件目录"]

    PluginLoader.PLUGINS_ROOT.mkdir(parents=True, exist_ok=True)
    installed = 0
    errors: List[str] = []

    for key in _BUILTIN_FORMAT_PLUGIN_KEYS:
        src = source_root / key
        manifest_path = src / "manifest.json"
        if not manifest_path.exists():
            errors.append(f"内置插件 {key} 缺少 manifest.json")
            continue

        try:
            manifest_data = json.loads(manifest_path.read_text(encoding="utf-8"))
            manifest = PluginManifest.model_validate(manifest_data)
            dst = PluginLoader.get_plugin_dir(key)
            shutil.copytree(src, dst, dirs_exist_ok=True, ignore=_ignore_runtime_artifacts)
            if key == "cc.foxel.officeviewer":
                route_path = dst / "backend" / "routes" / "office_preview.py"
                route_path.parent.mkdir(parents=True, exist_ok=True)
                route_path.write_text(_office_preview_route_source(), encoding="utf-8")
                frontend_path = dst / "frontend" / "index.js"
                frontend_path.parent.mkdir(parents=True, exist_ok=True)
                frontend_path.write_text(_office_frontend_source(), encoding="utf-8")
                manifest_data["description"] = "Office 文档本地查看器，支持 docx/xlsx/pptx 预览；不依赖公开分享链接。"
                manifest_data["backend"] = {
                    "routes": [
                        {
                            "module": "backend/routes/office_preview.py",
                            "prefix": "/api/plugins/office-viewer",
                            "tags": ["office-viewer"],
                        }
                    ]
                }
                (dst / "manifest.json").write_text(
                    json.dumps(manifest_data, ensure_ascii=False, indent=2),
                    encoding="utf-8",
                )
                manifest = PluginManifest.model_validate(manifest_data)

            rec = await Plugin.get_or_none(key=manifest.key)
            if rec is None:
                rec = Plugin(key=manifest.key)
            _manifest_to_record(rec, manifest)
            await rec.save()
            installed += 1
        except Exception as exc:
            logger.exception("内置插件 %s 导入失败", key)
            errors.append(f"内置插件 {key} 导入失败: {exc}")

    return installed, errors


async def load_installed_plugins(app: "FastAPI") -> Tuple[int, List[str]]:
    """
    加载所有已安装的插件

    Args:
        app: FastAPI 应用实例

    Returns:
        (成功加载数量, 错误列表)
    """
    from models.database import Plugin

    errors: List[str] = []
    loaded_count = 0

    try:
        plugins = await Plugin.all()
    except Exception as e:
        logger.error(f"查询插件列表失败: {e}")
        return 0, [f"查询插件列表失败: {e}"]

    for plugin in plugins:
        if not plugin.key:
            continue

        try:
            # 获取 manifest
            manifest = None
            if plugin.manifest:
                try:
                    manifest = PluginManifest.model_validate(plugin.manifest)
                except Exception:
                    # 尝试从文件系统读取
                    manifest = PluginLoader.read_manifest(plugin.key)
            else:
                manifest = PluginLoader.read_manifest(plugin.key)

            if not manifest:
                logger.warning(f"插件 {plugin.key} 缺少 manifest，跳过加载")
                continue

            # 加载后端路由
            loaded_routes: List[str] = []
            if manifest.backend and manifest.backend.routes:
                try:
                    routers = PluginLoader.load_all_routes(plugin.key, manifest)
                    for router in routers:
                        app.include_router(router)
                        loaded_routes.append(router.prefix)
                    logger.info(f"插件 {plugin.key} 加载了 {len(routers)} 个路由")
                except PluginLoadError as e:
                    errors.append(f"插件 {plugin.key} 路由加载失败: {e}")
                    logger.error(f"插件 {plugin.key} 路由加载失败: {e}")

            # 加载处理器
            loaded_processors: List[str] = []
            if manifest.backend and manifest.backend.processors:
                try:
                    processor_types = PluginLoader.load_all_processors(plugin.key, manifest)
                    loaded_processors = processor_types
                    logger.info(f"插件 {plugin.key} 注册了 {len(processor_types)} 个处理器")
                except PluginLoadError as e:
                    errors.append(f"插件 {plugin.key} 处理器加载失败: {e}")
                    logger.error(f"插件 {plugin.key} 处理器加载失败: {e}")

            # 更新数据库记录
            plugin.loaded_routes = loaded_routes if loaded_routes else None
            plugin.loaded_processors = loaded_processors if loaded_processors else None
            await plugin.save()

            loaded_count += 1
            logger.info(f"插件 {plugin.key} 加载完成")

        except Exception as e:
            error_msg = f"插件 {plugin.key} 加载异常: {e}"
            errors.append(error_msg)
            logger.exception(error_msg)

    return loaded_count, errors


async def init_plugins(app: "FastAPI") -> None:
    """
    初始化插件系统

    在应用启动时调用
    """
    logger.info("开始加载已安装插件...")

    bootstrap_count, bootstrap_errors = await bootstrap_builtin_format_plugins()
    if bootstrap_errors:
        for error in bootstrap_errors:
            logger.warning("  - %s", error)
    elif bootstrap_count:
        logger.info("已导入 %d 个内置格式插件", bootstrap_count)

    loaded_count, errors = await load_installed_plugins(app)

    if errors:
        logger.warning(f"插件加载完成，共 {loaded_count} 个成功，{len(errors)} 个错误")
        for error in errors:
            logger.warning(f"  - {error}")
    else:
        logger.info(f"插件加载完成，共 {loaded_count} 个插件")
